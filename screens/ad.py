from kivy.core.window import Window
from kivy.lang import Builder
from kivy.uix.label import Label
from kivy.uix.screenmanager import ScreenManager, Screen
from kivy.clock import Clock
from kivy.uix.video import Video
from DbFuncs import db
import base64
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.videoplayer import VideoPlayer
from kivy.uix.video import Video    
from kivy.uix.image import Image
from pptx import Presentation
from config import utils
import os
from io import BytesIO
import subprocess
from Observers.Observer import Observer
from config.utils import getFileLock, FILELOCK_LOCK

from dotenv import load_dotenv
load_dotenv(override=True)
# import dotenv

# dotenv_file = dotenv.find_dotenv()
# dotenv.load_dotenv(dotenv_file)

class AdScreen(Screen, Observer):
    def __init__(self, **kwargs):
        super(AdScreen, self).__init__(**kwargs)
        self.videoPlayerLayout = VideoPlayerLayout("")
        self.videoPlayerLayout.manager = self.manager
        self.videoPlayerLayout.bind(on_touch_up=self.on_touch_screen)
        self.videoPlayerLayout.id = "video_layout"
        self.add_widget(self.videoPlayerLayout)
        self.file_path = ''
        self.current_video_id = -1
        self.currentVideoPath = None

    def on_pre_enter(self, *args):
        self.update()

    def update(self):
        Clock.schedule_once(self.retrieve_layout, 0)

    def retrieve_layout(self, dt):
        ad = db.get_ad()
        if ad:
            if ad.type == 'MP4':
                adDirPath = utils.get_main_script_dir() + '\\ad'

                try:
                    self.videoPlayerLayout.player.unload()
                    #self.videoPlayerLayout.player.state = 'stop'
                    #self.videoPlayerLayout.player.source = None
                except:
                    pass

                if getFileLock(FILELOCK_LOCK).acquire(False) == False:
                    return
                
                print('getfilelock_____finished============================================================================')
                try:
                    self.file_path = utils.write_to_file(ad.content, adDirPath)
                except:
                    print('write_to_file_except')
                    pass
                getFileLock(FILELOCK_LOCK).release()

                try:
                    self.videoPlayerLayout.change_video(self.file_path)
                except Exception as e:
                    print("change_video_error", e)
                    pass
        
            elif ad.type == 'PPT':
                # Save the blob data as a temporary .ppt file
                pptFileName = 'temp.ppt'
                pptxFileName = 'temp.pptx'
                # self.temp_file = './' + pptFileName
                # self.temp_file = utils.get_main_script_dir() + os.environ.get('pptPath') + pptFileName

                # ......./img/
                dirPath = utils.get_main_script_dir() + os.environ.get('pptPath')

                utils.write_to_file(ad.content, dirPath + pptFileName)

                # Convert .ppt file into .pptx file
                # Define the command as a list of arguments
                command = ['soffice', '--headless', '--convert-to', 'pptx', '--outdir', 
                           dirPath, dirPath + pptFileName]
                
                # Execute the command
                subprocess.run(command)


                presentation = Presentation(dirPath + pptxFileName)

                slides = presentation.slides

                for slide in slides:
                    # Iterate through the shapes in the slide
                    for shape in slide.shapes:
                        # Check if the shape is an image
                        if shape.shape_type == 13:  # 13 corresponds to image shape type
                            # Extract the image
                            image = shape.image
                            image_bytes = image.blob
                            image_data = BytesIO(image_bytes)
                            
                            # Convert the BytesIO object to a base64-encoded string
                            base64_image = base64.b64encode(image_data.getvalue()).decode()
                            
                            # Create an Image widget with the base64-encoded string as the source
                            slide_image = Image(source=f"data:image/png;base64,{base64_image}", height=400)
                            slide_image.bind(on_touch_down=self.touch_screen)
                            self.add_widget(slide_image)
                
                os.remove(dirPath + pptFileName)
                os.remove(dirPath + pptxFileName)

        else:
            label_widget = Label(text='Ads not found', color=(1,0,0,1))
            label_widget.bind(on_touch_up=self.on_touch_screen)
            self.add_widget(label_widget)
        
    def on_touch_screen(self, instance, touch):
        if instance.collide_point(*touch.pos):
            self.manager.current = 'List'

class VideoPlayerLayout(BoxLayout):
    def __init__(self, temp_file, **kwargs):
        super(VideoPlayerLayout, self).__init__(**kwargs)
        self.manager = None
        self.file = None

        # Create a VideoPlayer widget
        self.player = Video(state='play', options={'eos':'loop'})

        # Add the VideoPlayer widget to the BoxLayout
        self.add_widget(self.player)
    
    def change_video(self, file):
        try:
            # self.player.state = 'stop'
            previewUrl = utils.get_main_script_dir() + '\\img\\preview.png'
            self.player.preview = previewUrl
            self.player.source = file
            self.player.state = 'play'
        except Exception as e:
            print('play_error', e)      
